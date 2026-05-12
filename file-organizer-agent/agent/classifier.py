"""
agent/classifier.py — content-aware classifier using a local Ollama LLM.

Pipeline:
  1. Extract a text snippet from the file (PDF, DOCX, TXT, …)
  2. Scan the University folder tree to get available semesters + lectures
  3. Ask Ollama to pick the best (semester, lecture) match
  4. Return a ClassifyResult with the target path
"""

from __future__ import annotations
import json
from dataclasses import dataclass
from pathlib import Path

import ollama
import fitz                     # pymupdf
import docx                     # python-docx
from pptx import Presentation   # python-pptx
import base64
from io import BytesIO
from PIL import Image 
from rich.console import Console

from config import (
    UNIVERSITY_DIR, OLLAMA_MODEL, OLLAMA_URL,
    MAX_CONTENT_CHARS, READABLE_EXTENSIONS, IMAGE_EXTENSIONS, CONFIDENCE_THRESHOLD,
    LECTURE_SUBFOLDERS,
)

console = Console()


# ---------------------------------------------------------------------------
# Result type
# ---------------------------------------------------------------------------
@dataclass
class ClassifyResult:
    semester:   str        # e.g. "Semester 4"
    lecture:    str        # e.g. "Operating Systems"
    subfolder:  str        # e.g. "Assignments"
    confidence: float
    reason:     str
    error:      str = ""

    @property
    def target_dir(self) -> Path:
        return UNIVERSITY_DIR / self.semester / self.lecture / self.subfolder

    @property
    def is_confident(self) -> bool:
        return self.confidence >= CONFIDENCE_THRESHOLD


# ---------------------------------------------------------------------------
# Folder scanner
# ---------------------------------------------------------------------------
def scan_university_tree() -> dict[str, list[str]]:
    """
    Returns {semester_name: [lecture_name, …]} by reading the real folder tree.
    Example: {"1. Semester": ["Mathematik 1", "Informatik"], …}
    """
    tree: dict[str, list[str]] = {}
    if not UNIVERSITY_DIR.exists():
        raise FileNotFoundError(
            f"University folder not found: {UNIVERSITY_DIR}\n"
            "Create it first or update UNIVERSITY_DIR in config.py"
        )
    for sem_dir in sorted(UNIVERSITY_DIR.iterdir()):
        if sem_dir.is_dir() and not sem_dir.name.startswith("."):
            lectures = [
                d.name for d in sorted(sem_dir.iterdir())
                if d.is_dir() and not d.name.startswith(".")
            ]
            tree[sem_dir.name] = lectures
    return tree


# ---------------------------------------------------------------------------
# Content extractor
# ---------------------------------------------------------------------------
def extract_text(filepath: Path) -> str:
    """Extract a short text snippet from the file for LLM context."""
    ext = filepath.suffix.lower()
    text = ""

    try:
        if ext == ".pdf":
            doc = fitz.open(str(filepath))
            for page in doc:
                text += page.get_text()
                if len(text) >= MAX_CONTENT_CHARS:
                    break
            doc.close()

        elif ext in {".docx", ".doc"}:
            document = docx.Document(str(filepath))
            text = "\n".join(p.text for p in document.paragraphs)
            
        elif ext == {".pptx", ".ppt"}:
            prs = Presentation(str(filepath))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text += para.text + "\n"
                    if len(text) >= MAX_CONTENT_CHARS:
                        break
                if len(text) >= MAX_CONTENT_CHARS:
                    break

        elif ext in READABLE_EXTENSIONS:
            text = filepath.read_text(encoding="utf-8", errors="ignore")

    except Exception as e:
        console.print(f"  [yellow]Content extract warning:[/yellow] {e}")

    return text[:MAX_CONTENT_CHARS].strip()

def load_image_base64(filepath: Path, max_px: int = 1024) -> str:
    """
    Load an image, downscale it so the longest edge ≤ max_px, and return
    a base64-encoded JPEG string ready for the Ollama vision API.
    Keeping images small is important — large images slow inference
    significantly and provide little extra information for classification.
    """
    with Image.open(filepath) as img:
        img = img.convert("RGB")          # drop alpha, normalise mode
        img.thumbnail((max_px, max_px))   # resize in-place, keeps aspect ratio
        buf = BytesIO()
        img.save(buf, format="JPEG", quality=85)
        return base64.b64encode(buf.getvalue()).decode("utf-8")

# ---------------------------------------------------------------------------
# Classifier
# ---------------------------------------------------------------------------
class FileClassifier:
    def __init__(self):
        # Validate Ollama is reachable at startup
        try:
            client = ollama.Client(host=OLLAMA_URL)
            client.list()   # raises if Ollama isn't running
            self.client = client
            self.model = OLLAMA_MODEL
        except Exception:
            raise RuntimeError(
                "Cannot reach Ollama at " + OLLAMA_URL + "\n"
            )
        self.tree = scan_university_tree()
        console.print(f"  Loaded {sum(len(v) for v in self.tree.values())} lectures "
                      f"across {len(self.tree)} semesters from University folder.")

    # ------------------------------------------------------------------
    def classify(self, filepath: Path) -> ClassifyResult:
        ext       = filepath.suffix.lower()
        content   = ""
        image_b64 = None

        if ext in IMAGE_EXTENSIONS:
            try:
                image_b64 = load_image_base64(filepath)
            except Exception as e:
                return self._fallback(f"Could not load image: {e}")
        else:
            content = extract_text(filepath)

        return self._llm_classify(filepath.name, content, image_b64)

    # ------------------------------------------------------------------
    def _llm_classify(
        self,
        filename:  str,
        content:   str       = "",
        image_b64: str | None = None,
    ) -> ClassifyResult:
        tree_str = json.dumps(self.tree, ensure_ascii=False, indent=2)

        content_section = (
            "(see attached image)"          if image_b64 else
            content                         if content   else
            "(no readable content — use filename only)"
        )
        
        prompt = f"""You are a file organiser for a student.
Given a filename and a content snippet, choose the best semester folder,
lecture subfolder, and file-type subfolder from the structures below.

Available semester/lecture structure (JSON):
{tree_str}

Available subfolders under every lecture:
{LECTURE_SUBFOLDERS}

Filename: {filename}
Content:
---
{content_section}
---

Subfolder rules:
- Use "Assignments" for exercises, homework, problem sets, solutions, or lab sheets.
- Use "Lecture Notes" for slides, scripts, summaries, or lecturer-provided material.
- When in doubt, use "Lecture Notes".

Respond ONLY with a JSON object — no markdown, no explanation outside it:
{{
  "semester":  "<exact semester folder name>",
  "lecture":   "<exact lecture folder name>",
  "subfolder": "<Assignments or Lecture Notes>",
  "confidence": <float 0.0-1.0>,
  "reason": "<one sentence, max 15 words>"
}}
"""
        message = (
            {
                "role":    "user",
                "content": prompt,
                "images":  [image_b64],   # ← separate key, content stays a plain string
            }
            if image_b64 else
            {"role": "user", "content": prompt}
        )

        try:
            response = self.client.chat(
                model=OLLAMA_MODEL,
                messages=[message],
                format="json",
                options={"temperature": 0},
            )
            raw = response["message"]["content"]
            data = json.loads(raw)

            semester = data.get("semester", "")
            lecture  = data.get("lecture",  "")
            subfolder = data.get("subfolder", "")

            # Validate against real folder tree
            if subfolder not in LECTURE_SUBFOLDERS:
                subfolder = "Lecture Notes"
            if semester not in self.tree:
                return self._fallback(f"Unknown semester '{semester}' returned by LLM")
            if lecture not in self.tree[semester]:
                return self._fallback(
                    f"Lecture '{lecture}' not in {semester}", semester=semester
                )

            result = ClassifyResult(
                semester=semester,
                lecture=lecture,
                subfolder=subfolder,
                confidence=float(data.get("confidence", 0.5)),
                reason=data.get("reason", "LLM classification."),
            )

            if not result.is_confident:
                console.print(
                    f"  [yellow]⚠ Low confidence ({result.confidence:.0%}) — "
                    f"double-check placement[/yellow]"
                )

            return result

        except json.JSONDecodeError as e:
            return self._fallback(f"JSON parse error: {e}")
        except Exception as e:
            return self._fallback(f"LLM call failed: {e}")

    def _fallback(self, reason: str, semester: str = "") -> ClassifyResult:
        """When classification fails, pick the first semester as a safe landing zone."""
        first_sem = semester or (list(self.tree.keys())[0] if self.tree else "")
        console.print(f"  [red]Fallback:[/red] {reason}")
        return ClassifyResult(
            semester=first_sem,
            lecture="_Unsorted",
            subfolder="Lecture Notes",
            confidence=0.0,
            reason=reason,
            error=reason,
        )

    def reload_tree(self) -> None:
        """Call this if you add new semester/lecture folders at runtime."""
        self.tree = scan_university_tree()
        console.print("  [dim]Folder tree reloaded.[/dim]")